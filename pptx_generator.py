"""
pptx_generator.py

Turns the CSV files produced by stats_processor.py (a team-comparison
KPI table + per-player stat tables) into a formatted PowerPoint
(.pptx) report styled after a StatsBomb-style match recap deck:

  1. Title slide - dark background, scoreline, team monogram
  2. Team Totals - one horizontal grouped-bar chart of the 12
     headline KPIs
  3. One slide per headline KPI - team-total callout + side-by-side
     player-by-player horizontal bar charts
  4. (optional) Supplementary KPI table (PPDA + secondary stats)
  5. (optional) Average-position image slide
  6. Full-Time summary - categories won / lost / tied

This module has no Streamlit dependency - it can be run standalone or
imported by streamlit_app.py.
"""

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

import stats_processor as sp


# ============================================================
# LAYOUT / DESIGN CONSTANTS
# ============================================================

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

BG_DARK = RGBColor(0x0E, 0x0E, 0x0E)
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xED, 0xE6, 0xD6)
GREY_LABEL = RGBColor(0x9A, 0x9A, 0x9A)
MUTED_TEXT = RGBColor(0x6B, 0x6B, 0x6B)
DARK_TEXT = RGBColor(0x18, 0x18, 0x18)
ROW_ALT_BG = RGBColor(0xF2, 0xF0, 0xEB)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

MAX_PLAYERS_PER_CHART = 16


# ============================================================
# COLOR HELPERS
# ============================================================

def _hex_to_rgb(hex_color):
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"'{hex_color}' is not a valid hex color (expected 6 hex digits)")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _readable_text_color(rgb_color):
    luminance = 0.299 * rgb_color[0] + 0.587 * rgb_color[1] + 0.114 * rgb_color[2]
    return DARK_TEXT if luminance > 150 else RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
# LOW-LEVEL SLIDE HELPERS
# ============================================================

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_background(slide, prs, rgb_color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, italic=False,
                  color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=BODY_FONT,
                  anchor=None, letter_spaced=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    display_text = " ".join(text) if letter_spaced else text
    run.text = display_text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def _add_footer(slide, prs, left_text, right_text, on_dark=False):
    color = GREY_LABEL if on_dark else MUTED_TEXT
    _add_textbox(slide, Inches(0.5), Inches(SLIDE_HEIGHT_IN - 0.45), Inches(6.0), Inches(0.35),
                 left_text, size=10, color=color)
    _add_textbox(slide, Inches(SLIDE_WIDTH_IN - 6.5), Inches(SLIDE_HEIGHT_IN - 0.45), Inches(6.0), Inches(0.35),
                 right_text, size=10, color=color, align=PP_ALIGN.RIGHT)


def _add_monogram(slide, center_x, center_y, diameter, letter, fill_color, text_color):
    left = Emu(int(center_x - diameter / 2))
    top = Emu(int(center_y - diameter / 2))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Emu(int(diameter)), Emu(int(diameter)))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = letter
    run.font.size = Pt(int(diameter / Inches(1) * 30))
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = HEADER_FONT
    return circle


def _add_decorative_arc(slide, prs):
    """A thin quarter-circle outline peeking in from the bottom-left corner."""
    diameter = Inches(3.2)
    left = Inches(-1.6)
    top = Inches(SLIDE_HEIGHT_IN - 1.6)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    circle.line.width = Pt(1)
    circle.shadow.inherit = False
    return circle


def _add_score_pill(slide, left, top, score, on_color, text_color, diameter=Inches(0.62)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = on_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(score)
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = HEADER_FONT
    return circle


def _style_table(table, header_color, header_text_color):
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = header_text_color
                r.font.size = Pt(14)
                r.font.name = BODY_FONT

    for row_idx in range(1, len(table.rows)):
        bg = ROW_ALT_BG if row_idx % 2 == 0 else BG_LIGHT
        for cell in table.rows[row_idx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = DARK_TEXT
                    r.font.name = BODY_FONT


def _style_single_series_bar_chart(chart, rgb_color, reverse_categories=True):
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.color.rgb = DARK_TEXT
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = rgb_color
    series.format.line.fill.background()

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.format.line.color.rgb = RGBColor(0xD5, 0xD5, 0xD5)
    category_axis.reverse_order = reverse_categories

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.minimum_scale = 0


# ============================================================
# 1. TITLE SLIDE
# ============================================================

def add_title_slide(prs, team_name, team_color, opponent_name, opponent_color,
                     team_score=None, opponent_score=None,
                     kicker="STATSBOMB MATCH RECAP",
                     subtitle="STATISTICAL MATCH ANALYSIS  \u2022  PLAYER & TEAM PERFORMANCE BREAKDOWN",
                     date_label=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_DARK)
    _add_decorative_arc(slide, prs)

    _add_textbox(slide, Inches(0.8), Inches(1.55), Inches(9.0), Inches(0.4),
                 kicker, size=13, color=GREY_LABEL, letter_spaced=True)

    _add_textbox(slide, Inches(0.75), Inches(2.0), Inches(9.0), Inches(1.0),
                 team_name.upper(), size=44, bold=True, color=team_color, font_name=HEADER_FONT)
    if team_score is not None:
        _add_score_pill(slide, Inches(0.8 + 0.11 * len(team_name)), Inches(1.95), team_score,
                         RGBColor(0xFF, 0xFF, 0xFF), DARK_TEXT)

    _add_textbox(slide, Inches(0.8), Inches(2.95), Inches(9.0), Inches(0.7),
                 opponent_name.upper(), size=28, italic=True, color=CREAM, font_name=HEADER_FONT)

    _add_textbox(slide, Inches(0.8), Inches(3.85), Inches(9.5), Inches(0.4),
                 subtitle, size=12, color=GREY_LABEL, letter_spaced=False)

    monogram_letter = team_name.strip()[0].upper() if team_name.strip() else "?"
    _add_monogram(slide, Inches(11.1), Inches(3.7), Inches(2.6), monogram_letter, CREAM, BG_DARK)

    footer_right = f"{team_name.upper()}   vs   {opponent_name.upper()}"
    _add_footer(slide, prs, date_label or "", footer_right, on_dark=True)

    return slide


# ============================================================
# 2. TEAM TOTALS SLIDE
# ============================================================

def add_team_totals_slide(prs, comparison_df, team_name, opponent_name, team_color, opponent_color,
                           metrics=None):
    metrics = metrics or sp.MAIN_METRICS
    labels = [label for label, _ in metrics]

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(6.0), Inches(0.55),
                 "TEAM TOTALS", size=26, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    legend_x = Inches(7.6)
    chip1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, Inches(0.5), Inches(0.28), Inches(0.28))
    chip1.fill.solid()
    chip1.fill.fore_color.rgb = team_color
    chip1.line.fill.background()
    chip1.shadow.inherit = False
    _add_textbox(slide, legend_x + Inches(0.38), Inches(0.44), Inches(2.3), Inches(0.4),
                 team_name, size=13, bold=True, color=DARK_TEXT)

    chip2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x + Inches(2.6), Inches(0.5), Inches(0.28), Inches(0.28))
    chip2.fill.solid()
    chip2.fill.fore_color.rgb = opponent_color
    chip2.line.fill.background()
    chip2.shadow.inherit = False
    _add_textbox(slide, legend_x + Inches(2.98), Inches(0.44), Inches(2.5), Inches(0.4),
                 opponent_name, size=13, bold=True, color=DARK_TEXT)

    rows = comparison_df.set_index("Metric").reindex(labels).reset_index()

    chart_data = CategoryChartData()
    chart_data.categories = list(reversed(rows["Metric"].tolist()))
    chart_data.add_series(team_name, list(reversed(rows[team_name].fillna(0).tolist())))
    chart_data.add_series(opponent_name, list(reversed(rows[opponent_name].fillna(0).tolist())))

    x, y, cx, cy = Inches(0.6), Inches(1.1), Inches(SLIDE_WIDTH_IN - 1.2), Inches(5.9)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.gap_width = 60
    plot.overlap = -10

    series = plot.series
    series[0].format.fill.solid()
    series[0].format.fill.fore_color.rgb = team_color
    series[0].format.line.fill.background()
    series[1].format.fill.solid()
    series[1].format.fill.fore_color.rgb = opponent_color
    series[1].format.line.fill.background()

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(11)
    category_axis.format.line.color.rgb = RGBColor(0xD5, 0xD5, 0xD5)

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.minimum_scale = 0

    _add_footer(slide, prs, "Team Totals \u2014 All Statistics", f"{team_name.upper()} vs {opponent_name.upper()}")
    return slide


# ============================================================
# 3. PER-METRIC PLAYER BREAKDOWN SLIDES
# ============================================================

def add_metric_slide(prs, label, player_col, team_value, opponent_value,
                      team_players_df, opponent_players_df,
                      team_name, opponent_name, team_color, opponent_color,
                      stat_index=None, stat_total=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.55),
                 label.upper(), size=24, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)
    if stat_index is not None and stat_total is not None:
        _add_textbox(slide, Inches(0.5), Inches(0.78), Inches(4.0), Inches(0.3),
                     f"Statistic {stat_index} of {stat_total}", size=11, color=MUTED_TEXT)

    _add_textbox(slide, Inches(9.2), Inches(0.32), Inches(3.6), Inches(0.3),
                 "TEAM TOTAL", size=11, color=MUTED_TEXT, align=PP_ALIGN.RIGHT)

    def _fmt(v):
        if v is None or pd.isna(v):
            return "-"
        return str(int(v)) if float(v).is_integer() else f"{v:.1f}"

    _add_textbox(slide, Inches(9.0), Inches(0.55), Inches(1.8), Inches(0.75),
                 _fmt(team_value), size=32, bold=True, color=team_color, align=PP_ALIGN.CENTER,
                 font_name=HEADER_FONT)
    _add_textbox(slide, Inches(10.8), Inches(0.55), Inches(1.9), Inches(0.75),
                 _fmt(opponent_value), size=32, bold=True, color=opponent_color, align=PP_ALIGN.CENTER,
                 font_name=HEADER_FONT)
    _add_textbox(slide, Inches(9.0), Inches(1.25), Inches(1.8), Inches(0.3),
                 team_name, size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(10.8), Inches(1.25), Inches(1.9), Inches(0.3),
                 opponent_name, size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    half_w = Inches((SLIDE_WIDTH_IN - 1.4) / 2)
    chart_top = Inches(2.0)
    chart_h = Inches(4.6)
    left_left = Inches(0.6)
    right_left = Inches(0.6) + half_w + Inches(0.2)

    for side_left, df, name, color in (
        (left_left, team_players_df, team_name, team_color),
        (right_left, opponent_players_df, opponent_name, opponent_color),
    ):
        _add_textbox(slide, side_left, chart_top - Inches(0.35), half_w, Inches(0.3),
                     name.upper(), size=12, bold=True, color=DARK_TEXT)

        if df is None or player_col not in df.columns or "Player" not in df.columns or df.empty:
            _add_textbox(slide, side_left, chart_top, half_w, Inches(0.4),
                         "No player data available", size=11, color=MUTED_TEXT)
            continue

        sub = df[["Player", player_col]].fillna(0).copy()
        sub = sub.sort_values(player_col, ascending=False).head(MAX_PLAYERS_PER_CHART)

        chart_data = CategoryChartData()
        chart_data.categories = sub["Player"].tolist()
        chart_data.add_series(name, sub[player_col].tolist())

        gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, side_left, chart_top, half_w, chart_h, chart_data)
        _style_single_series_bar_chart(gframe.chart, color)

    footer_left = label
    footer_right = f"{team_name.upper()} vs {opponent_name.upper()}"
    _add_footer(slide, prs, footer_left, footer_right)
    return slide


# ============================================================
# 4. SUPPLEMENTARY KPI TABLE
# ============================================================

def add_supplementary_table_slide(prs, comparison_df, team_name, opponent_name, team_color, opponent_color,
                                   metrics=None, title="Additional KPIs"):
    metrics = metrics or ([("PPDA", "PPDA")] + sp.SUPPLEMENTARY_METRICS)
    labels = [label for label, _ in metrics]
    rows = comparison_df[comparison_df["Metric"].isin(labels)]
    rows = rows.set_index("Metric").reindex(labels).reset_index()

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)
    _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.6),
                 title, size=26, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    n_rows = len(rows) + 1
    table_top = Inches(1.4)
    table_left = Inches(2.7)
    table_width = Inches(SLIDE_WIDTH_IN - 5.4)
    table_height = Inches(min(4.8, 0.6 * n_rows))

    gshape = slide.shapes.add_table(n_rows, 3, table_left, table_top, table_width, table_height)
    table = gshape.table
    table.columns[0].width = Emu(int(table_width * 0.42))
    table.columns[1].width = Emu(int(table_width * 0.29))
    table.columns[2].width = Emu(int(table_width * 0.29))

    table.cell(0, 0).text = "KPI"
    table.cell(0, 1).text = team_name
    table.cell(0, 2).text = opponent_name

    for i, row in enumerate(rows.itertuples(index=False), start=1):
        table.cell(i, 0).text = str(row.Metric)
        val_team = getattr(row, team_name, row[1])
        val_opp = getattr(row, opponent_name, row[2])
        table.cell(i, 1).text = "-" if pd.isna(val_team) else str(val_team)
        table.cell(i, 2).text = "-" if pd.isna(val_opp) else str(val_opp)
        for c in range(3):
            table.cell(i, c).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    _style_table(table, team_color, _readable_text_color(team_color))
    _add_footer(slide, prs, "Supplementary Statistics", f"{team_name.upper()} vs {opponent_name.upper()}")
    return slide


# ============================================================
# 5. AVERAGE POSITION IMAGE SLIDE
# ============================================================

def add_image_slide(prs, title, image_path, team_name, team_color, caption=None):
    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_LIGHT)
    _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.6),
                 title, size=26, bold=True, color=DARK_TEXT, font_name=HEADER_FONT)

    max_w = Inches(SLIDE_WIDTH_IN - 1.6)
    max_h = Inches(5.1)
    top = Inches(1.2)

    from PIL import Image
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    aspect = img_w / img_h

    if (max_w / max_h) > aspect:
        height = max_h
        width = Emu(int(height * aspect))
    else:
        width = max_w
        height = Emu(int(width / aspect))

    left = Emu(int((prs.slide_width - width) / 2))
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    if caption:
        _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.4),
                     caption, size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    _add_footer(slide, prs, title, team_name.upper())
    return slide


# ============================================================
# 6. FULL-TIME SUMMARY SLIDE
# ============================================================

def add_fulltime_slide(prs, comparison_df, team_name, opponent_name, team_color, opponent_color,
                        metrics=None):
    metrics = metrics or sp.MAIN_METRICS
    labels = [label for label, _ in metrics]
    rows = comparison_df[comparison_df["Metric"].isin(labels)]

    team_wins = int((rows[team_name] > rows[opponent_name]).sum())
    opponent_wins = int((rows[opponent_name] > rows[team_name]).sum())
    ties = int(len(rows) - team_wins - opponent_wins)

    slide = _blank_slide(prs)
    _add_background(slide, prs, BG_DARK)
    _add_decorative_arc(slide, prs)

    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(6.0), Inches(0.4),
                 "F U L L - T I M E", size=14, color=GREY_LABEL)

    if team_wins > opponent_wins:
        headline = f"{team_name} leads the match statistics"
    elif opponent_wins > team_wins:
        headline = f"{opponent_name} leads the match statistics"
    else:
        headline = "Match statistics are evenly split"

    _add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.7),
                 headline, size=30, bold=True, color=CREAM, font_name=HEADER_FONT)

    # Ensure opponent_color is legible on BG_DARK if black is selected
    opp_luminance = 0.299 * opponent_color[0] + 0.587 * opponent_color[1] + 0.114 * opponent_color[2]
    opp_display_color = CREAM if opp_luminance < 60 else opponent_color

    stat_specs = [
        (team_wins, f"categories won\nby {team_name}", team_color),
        (opponent_wins, f"categories won\nby {opponent_name}", opp_display_color),
        (ties, "categories tied", GREY_LABEL),
    ]

    box_w = Inches(3.4)
    gap = Inches(0.6)
    total_w = box_w * 3 + gap * 2
    start_x = Emu(int((prs.slide_width - total_w) / 2))
    top = Inches(2.9)

    for i, (value, caption, color) in enumerate(stat_specs):
        left = Emu(int(start_x + i * (box_w + gap)))
        _add_textbox(slide, left, top, box_w, Inches(1.3),
                     str(value), size=64, bold=True, color=color, align=PP_ALIGN.CENTER,
                     font_name=HEADER_FONT)
        cap_box = slide.shapes.add_textbox(left, top + Inches(1.35), box_w, Inches(0.7))
        tf = cap_box.text_frame
        tf.word_wrap = True
        for j, line in enumerate(caption.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = GREY_LABEL
            run.font.name = BODY_FONT

    _add_footer(slide, prs, "Match Recap Complete", f"{team_name.upper()} vs {opponent_name.upper()}", on_dark=True)
    return slide


# ============================================================
# MAIN ENTRY POINT
# ============================================================

def generate_report(
    comparison_csv,
    team_name,
    team_color,
    opponent_name,
    opponent_color,
    output_path,
    team_score=None,
    opponent_score=None,
    date_label=None,
    kicker="STATSBOMB MATCH RECAP",
    subtitle="STATISTICAL MATCH ANALYSIS  \u2022  PLAYER & TEAM PERFORMANCE BREAKDOWN",
    team_player_stats_csv=None,
    opponent_player_stats_csv=None,
    avg_position_image=None,
    include_supplementary_table=True,
):
    team_rgb = _hex_to_rgb(team_color)
    opponent_rgb = _hex_to_rgb(opponent_color)

    comparison_df = pd.read_csv(comparison_csv)
    for col in (team_name, opponent_name):
        comparison_df[col] = pd.to_numeric(comparison_df[col], errors="coerce")

    team_players_df = pd.read_csv(team_player_stats_csv) if team_player_stats_csv is not None else None
    opponent_players_df = pd.read_csv(opponent_player_stats_csv) if opponent_player_stats_csv is not None else None

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # 1. Title slide
    add_title_slide(prs, team_name, team_rgb, opponent_name, opponent_rgb,
                     team_score=team_score, opponent_score=opponent_score,
                     kicker=kicker, subtitle=subtitle, date_label=date_label)

    # 2. Team totals slide
    add_team_totals_slide(prs, comparison_df, team_name, opponent_name, team_rgb, opponent_rgb)

    # 3. Per-metric player-breakdown slides
    main_rows = comparison_df.set_index("Metric").reindex([l for l, _ in sp.MAIN_METRICS])
    stat_total = len(sp.MAIN_METRICS)
    for i, (label, player_col) in enumerate(sp.MAIN_METRICS, start=1):
        team_value = main_rows.loc[label, team_name] if label in main_rows.index else None
        opponent_value = main_rows.loc[label, opponent_name] if label in main_rows.index else None
        add_metric_slide(
            prs, label, player_col, team_value, opponent_value,
            team_players_df, opponent_players_df,
            team_name, opponent_name, team_rgb, opponent_rgb,
            stat_index=i, stat_total=stat_total,
        )

    # 4. Supplementary KPI table (PPDA + secondary stats)
    if include_supplementary_table:
        add_supplementary_table_slide(prs, comparison_df, team_name, opponent_name, team_rgb, opponent_rgb)

    # 5. Average position image slide (optional)
    if avg_position_image is not None:
        add_image_slide(prs, f"{team_name} - Average Position", avg_position_image, team_name, team_rgb)

    # 6. Full-time summary slide
    add_fulltime_slide(prs, comparison_df, team_name, opponent_name, team_rgb, opponent_rgb)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)
